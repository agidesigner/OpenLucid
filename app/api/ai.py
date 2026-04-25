from __future__ import annotations

import io
import logging

import json as json_mod

import httpx
from fastapi import APIRouter, Depends, File, Form, HTTPException, UploadFile
from fastapi.responses import StreamingResponse
from sqlalchemy.ext.asyncio import AsyncSession

from app.adapters.ai import get_ai_adapter, OpenAICompatibleAdapter, StubAIAdapter
from app.api.deps import get_db
from app.schemas.ai import (
    ExtractTextResponse,
    InferOfferKnowledgeRequest,
    InferOfferKnowledgeResponse,
    InferredKnowledgeItem,
)

logger = logging.getLogger(__name__)

router = APIRouter(prefix="/ai", tags=["ai"])


@router.get("/status")
async def ai_status(db: AsyncSession = Depends(get_db)):
    """Check whether a real LLM is configured and ready."""
    adapter = await get_ai_adapter(db, scene_key="knowledge")
    ready = not isinstance(adapter, StubAIAdapter)
    info = {}
    if isinstance(adapter, OpenAICompatibleAdapter):
        info = {"provider": adapter.provider, "model": adapter.model}
    return {"ready": ready, **info}


def _build_offer_data(body: InferOfferKnowledgeRequest) -> dict:
    """Build offer_data dict from request. Shared by stream and non-stream endpoints."""
    knowledge_items = []
    if body.existing_knowledge:
        knowledge_items = [
            {"knowledge_type": k.knowledge_type, "title": k.title, "content_raw": k.content_raw}
            for k in body.existing_knowledge
        ]
    return {
        "offer": {"name": body.name, "offer_type": body.offer_type, "description": body.description},
        "selling_points": [], "target_audiences": [], "target_scenarios": [],
        "knowledge_items": knowledge_items,
    }


def _infer_language_from_body(body: InferOfferKnowledgeRequest) -> str:
    """KB-centric language for knowledge inference.

    AI smart update / create-wizard KB generation has no UI language
    picker — the only signal is the brief text and any existing KB. Under
    the system-wide rule (presence of ``language`` is the explicit API
    override; anything else follows the content), we intentionally pass
    ``None`` for the explicit side so a Chinese UI uploading an English
    brief still gets English KB back. ``body.language`` is just the UI
    locale, not a user pick, so we don't forward it as authoritative.
    """
    from app.libs.lang_detect import resolve_output_language

    sample_parts: list[str] = [body.name or "", body.description or ""]
    if body.existing_knowledge:
        for k in body.existing_knowledge[:15]:
            sample_parts.append((k.title or "") + " " + (k.content_raw or "")[:500])
    sample = "\n".join(p for p in sample_parts if p)
    return resolve_output_language(None, sample, caller="infer_knowledge")


def _friendly_llm_error(e: Exception, adapter) -> str:
    """Turn an OpenAI SDK exception into a message that tells the user (a) which
    model failed, (b) what class of failure it was, (c) where to go next. The
    raw str(e) is often just "Request timed out." which sends users hunting."""
    from openai import APITimeoutError, APIConnectionError, RateLimitError, AuthenticationError, BadRequestError

    provider = getattr(adapter, "provider", "?")
    model = getattr(adapter, "model", "?")
    label = f"{provider}/{model}"

    if isinstance(e, APITimeoutError):
        return (
            f"模型 {label} 连续超时 3 次未响应。"
            f"请检查：(1) Settings → 模型配置中该 LLM 的 API 地址是否可达；"
            f"(2) 若走代理需确认代理服务正常；(3) 可在 Settings 中将 knowledge 场景切换到其他可用模型。"
        )
    if isinstance(e, APIConnectionError):
        return (
            f"模型 {label} 连接失败：{e}。"
            f"请检查 Settings → 模型配置中的 base_url 与网络连通性；或切换到其他模型。"
        )
    if isinstance(e, RateLimitError):
        return (
            f"模型 {label} 触发限流 / 额度耗尽：{e}。"
            f"请补充该 LLM 账户额度，或在 Settings 中切换到其他可用模型。"
        )
    if isinstance(e, AuthenticationError):
        return (
            f"模型 {label} 鉴权失败：{e}。"
            f"请在 Settings → 模型配置中核对 API Key。"
        )
    if isinstance(e, BadRequestError):
        return f"模型 {label} 拒绝请求：{e}。可能是 prompt 超长或参数不兼容，请反馈。"
    # Unknown — keep raw message but prefix with model label for traceability
    return f"模型 {label} 调用失败：{e}"


def _build_suggestions(raw: dict) -> dict[str, list[InferredKnowledgeItem]]:
    suggestions = {}
    for category, items in raw.items():
        suggestions[category] = [
            InferredKnowledgeItem(
                knowledge_type=category,
                title=item.get("title", ""),
                content_raw=item.get("content_raw", ""),
                confidence=item.get("confidence", 0.0),
            )
            for item in items
        ]
    return suggestions


@router.post("/infer-offer-knowledge", response_model=InferOfferKnowledgeResponse)
async def infer_offer_knowledge(body: InferOfferKnowledgeRequest, db: AsyncSession = Depends(get_db)):
    """Infer knowledge from offer info. Works for both creation (no existing) and update (with existing)."""
    offer_data = _build_offer_data(body)
    effective_lang = _infer_language_from_body(body)

    adapter = await get_ai_adapter(db, scene_key="knowledge")
    if isinstance(adapter, StubAIAdapter):
        raise HTTPException(
            status_code=503,
            detail="NO_LLM_CONFIGURED",
        )
    try:
        raw = await adapter.infer_knowledge(offer_data, effective_lang)
    except Exception as e:
        desc_len = len(body.description or "")
        logger.error(
            "infer-offer-knowledge failed | name=%s type=%s desc_len=%d lang=%s | %s",
            body.name, body.offer_type, desc_len, effective_lang, e,
            exc_info=True,
        )
        raise HTTPException(status_code=502, detail=_friendly_llm_error(e, adapter))

    if not raw:
        logger.error("infer-offer-knowledge returned empty | name=%s type=%s desc_len=%d lang=%s",
                      body.name, body.offer_type, len(body.description or ""), effective_lang)
        raise HTTPException(status_code=502, detail="LLM returned empty response")

    # Extract description generated by LLM (if present), remove from raw before building suggestions
    description = raw.pop("description", None) if isinstance(raw, dict) else None

    return InferOfferKnowledgeResponse(
        offer_name=body.name,
        description=description,
        suggestions=_build_suggestions(raw),
    )


@router.post("/infer-offer-knowledge-stream")
async def infer_offer_knowledge_stream(body: InferOfferKnowledgeRequest, db: AsyncSession = Depends(get_db)):
    """Streaming version: sends SSE events with thinking text, then final result."""
    offer_data = _build_offer_data(body)
    effective_lang = _infer_language_from_body(body)

    adapter = await get_ai_adapter(db, scene_key="knowledge")
    if isinstance(adapter, StubAIAdapter):
        raise HTTPException(status_code=503, detail="NO_LLM_CONFIGURED")
    if not isinstance(adapter, OpenAICompatibleAdapter):
        raise HTTPException(status_code=501, detail="Streaming not supported for this adapter")

    async def event_stream():
        try:
            async for event_type, data in adapter.infer_knowledge_stream(offer_data, effective_lang):
                if event_type == "thinking":
                    yield f"data: {json_mod.dumps({'type': 'thinking', 'text': data}, ensure_ascii=False)}\n\n"
                elif event_type == "error":
                    yield f"data: {json_mod.dumps({'type': 'error', 'detail': data}, ensure_ascii=False)}\n\n"
                elif event_type == "result":
                    description = data.pop("description", None) if isinstance(data, dict) else None
                    suggestions = {
                        cat: [item.model_dump() for item in items]
                        for cat, items in _build_suggestions(data).items()
                    }
                    result = {
                        "type": "result",
                        "offer_name": body.name,
                        "description": description,
                        "suggestions": suggestions,
                    }
                    yield f"data: {json_mod.dumps(result, ensure_ascii=False)}\n\n"
        except Exception as e:
            logger.error("infer-offer-knowledge-stream failed | name=%s | %s", body.name, e, exc_info=True)
            yield f"data: {json_mod.dumps({'type': 'error', 'detail': _friendly_llm_error(e, adapter)}, ensure_ascii=False)}\n\n"

    return StreamingResponse(event_stream(), media_type="text/event-stream")


def _extract_pdf_text(content: bytes) -> str:
    from PyPDF2 import PdfReader
    reader = PdfReader(io.BytesIO(content))
    parts = []
    for page in reader.pages:
        text = page.extract_text()
        if text:
            parts.append(text)
    return "\n".join(parts)


def _strip_jina_metadata(text: str) -> str:
    """Strip Jina Reader metadata headers. Leave content as-is — LLM handles markdown fine."""
    import re
    text = re.sub(r"^(Title|URL Source|Published Time|Markdown Content):.*$", "", text, flags=re.MULTILINE)
    text = re.sub(r"\n{3,}", "\n\n", text)
    return text.strip()


def _harvest_structured_metadata(raw_html: str) -> str:
    """Pull JSON-LD product/SoftwareApplication blocks + nav anchor text out
    of raw HTML before any stripping passes destroy them.

    Why this matters: product sites embed an authoritative "what this is"
    summary in ``<script type="application/ld+json">`` (schema.org
    SoftwareApplication / Product / Service). It enumerates the actual
    feature list, price, OS support, and category — exactly the fields
    our infer-knowledge prompt needs but normally has to guess from
    marketing copy. The previous ``_extract_url_text`` flow killed this
    in two places (Jina Reader's plain-text view drops scripts; the
    direct-fetch fallback explicitly strips ``<script>``), so an LLM
    looking at PrivacyCrop's homepage saw "Privacy Blur" hero copy and
    confidently summarised the whole product as "screenshot redaction"
    — missing the 6 other features the JSON-LD listed by name.

    Returns a labelled prefix block (or empty string if nothing useful
    was found). Callers prepend this to whatever the page-text extractor
    yields, so the LLM sees the structured metadata first.
    """
    import json as _json
    import re as _re
    if not raw_html:
        return ""

    sections: list[str] = []

    # 1) JSON-LD product-ish blocks — schema.org gives canonical product data
    try:
        ld_pattern = _re.compile(
            r'<script[^>]*type=["\']application/ld\+json["\'][^>]*>([\s\S]*?)</script>',
            _re.IGNORECASE,
        )
        relevant_types = {"SoftwareApplication", "Product", "Service", "WebApplication", "MobileApplication"}
        for m in ld_pattern.finditer(raw_html):
            blob = m.group(1).strip()
            if not blob:
                continue
            try:
                data = _json.loads(blob)
            except _json.JSONDecodeError:
                continue
            # @graph wraps multi-entity payloads
            entities = data.get("@graph") if isinstance(data, dict) and "@graph" in data else (
                data if isinstance(data, list) else [data]
            )
            for ent in entities or []:
                if not isinstance(ent, dict):
                    continue
                t = ent.get("@type")
                if isinstance(t, list):
                    matched = any(x in relevant_types for x in t)
                else:
                    matched = t in relevant_types
                if not matched:
                    continue
                # Trim to the high-signal fields the LLM actually needs.
                trimmed = {
                    k: ent.get(k) for k in
                    ("@type", "name", "description", "applicationCategory",
                     "operatingSystem", "browserRequirements", "featureList",
                     "offers", "category")
                    if ent.get(k) is not None
                }
                sections.append(_json.dumps(trimmed, ensure_ascii=False, indent=2))
    except Exception:
        pass

    # 2) Top-level <nav> anchor text — the sidebar / header navigation
    # usually IS the product's feature list, written by the site owner.
    # Pulling these as a flat list complements the JSON-LD when the
    # latter is missing or sparse.
    try:
        nav_pattern = _re.compile(r"<nav[^>]*>([\s\S]*?)</nav>", _re.IGNORECASE)
        anchor_pattern = _re.compile(r"<a[^>]*>([\s\S]*?)</a>", _re.IGNORECASE)
        tag_pattern = _re.compile(r"<[^>]+>")
        nav_items: list[str] = []
        # Cap at first 3 <nav> blocks; large sites have multi-level nav.
        nav_blocks = list(nav_pattern.finditer(raw_html))[:3]
        for nav_match in nav_blocks:
            for a_match in anchor_pattern.finditer(nav_match.group(1)):
                text = tag_pattern.sub(" ", a_match.group(1))
                text = " ".join(text.split())
                if text and 1 < len(text) <= 60 and text not in nav_items:
                    nav_items.append(text)
        if nav_items:
            sections.append("Navigation items:\n  - " + "\n  - ".join(nav_items[:30]))
    except Exception:
        pass

    if not sections:
        return ""
    body = "\n\n".join(sections)
    return f"[STRUCTURED PRODUCT METADATA — extracted from page schema/nav, treat as authoritative]\n{body}\n\n[/STRUCTURED PRODUCT METADATA]\n\n"


async def _extract_url_text(url: str) -> str:
    """Extract text from a URL using Jina Reader (handles JS-rendered SPA pages).

    Now also harvests JSON-LD ``SoftwareApplication`` / ``Product`` blocks
    and nav-anchor text from the raw HTML and prepends them as a labelled
    "[STRUCTURED PRODUCT METADATA]" prefix. This survives even when Jina
    Reader's plain-text view drops the original script tags. Without it
    the LLM only sees marketing hero copy and confidently mis-classifies
    multi-feature products by their loudest hero section.
    """
    import re

    # Fetch raw HTML once up front so we can both harvest structured
    # metadata AND fall back to it if Jina is unavailable.
    raw_html = ""
    try:
        async with httpx.AsyncClient(follow_redirects=True, timeout=15) as client:
            resp = await client.get(url, headers={"User-Agent": "Mozilla/5.0 (compatible; OpenLucid)"})
            resp.raise_for_status()
            raw_html = resp.text
    except Exception:
        raw_html = ""
    metadata_prefix = _harvest_structured_metadata(raw_html)

    jina_url = f"https://r.jina.ai/{url}"
    jina_error = None
    try:
        async with httpx.AsyncClient(follow_redirects=True, timeout=45) as client:
            resp = await client.get(jina_url, headers={"Accept": "text/plain"})
            resp.raise_for_status()
            text = _strip_jina_metadata(resp.text)
            if text:
                return metadata_prefix + text
    except Exception as e:
        jina_error = str(e)
        logger.warning("Jina Reader failed for %s: %s, falling back to direct fetch", url, jina_error)

    # Fallback: direct fetch + strip tags. Reuse the raw HTML we already
    # pulled to avoid a second network round-trip.
    if not raw_html:
        try:
            async with httpx.AsyncClient(follow_redirects=True, timeout=15) as client:
                resp = await client.get(url)
                resp.raise_for_status()
                raw_html = resp.content.decode("utf-8", errors="ignore")
        except Exception as e:
            raise HTTPException(400, f"Unable to access URL: {e}")

    text = re.sub(r"<script[^>]*>[\s\S]*?</script>", "", raw_html)
    text = re.sub(r"<style[^>]*>[\s\S]*?</style>", "", text)
    text = re.sub(r"<[^>]+>", " ", text)
    text = re.sub(r"\s+", " ", text).strip()

    # If extracted text is too short, the page likely requires JS rendering
    # and direct fetch only got an empty shell. Fail with a clear message.
    if len(text) < 200:
        raise HTTPException(
            400,
            "This page returned very little text content — it likely requires "
            "JavaScript to render. Please copy and paste the page content directly.",
        )

    # Detect code/script noise in extracted text (SPA frameworks leak inline JS)
    noise_chars = sum(1 for c in text if c in '{}();=>')
    noise_ratio = noise_chars / len(text) if text else 0
    if noise_ratio > 0.03:
        logger.warning("URL text has high code noise (%.1f%%): %s", noise_ratio * 100, url)
        raise HTTPException(
            400,
            "The extracted page content contains too much script/code noise — "
            "this site likely requires JavaScript to render properly. "
            "Please copy and paste the page content directly.",
        )

    return metadata_prefix + text


def _extract_docx_text(content: bytes) -> str:
    from docx import Document
    doc = Document(io.BytesIO(content))
    return "\n".join(p.text for p in doc.paragraphs if p.text.strip())


def _extract_pptx_text(content: bytes) -> str:
    from pptx import Presentation
    prs = Presentation(io.BytesIO(content))
    lines = []
    for i, slide in enumerate(prs.slides, 1):
        slide_lines = []
        for shape in slide.shapes:
            if shape.has_text_frame:
                for para in shape.text_frame.paragraphs:
                    t = "".join(r.text for r in para.runs).strip()
                    if t:
                        slide_lines.append(t)
            elif shape.has_table:
                for row in shape.table.rows:
                    cells = [c.text.strip() for c in row.cells if c.text.strip()]
                    if cells:
                        slide_lines.append(" | ".join(cells))
        if slide_lines:
            lines.append(f"[Slide {i}]")
            lines.extend(slide_lines)
    return "\n".join(lines)


def _extract_excel_text(content: bytes) -> str:
    from openpyxl import load_workbook
    wb = load_workbook(io.BytesIO(content), read_only=True, data_only=True)
    lines = []
    for sheet in wb.worksheets:
        if len(wb.worksheets) > 1:
            lines.append(f"[{sheet.title}]")
        for row in sheet.iter_rows(values_only=True):
            cells = [str(c).strip() for c in row if c is not None]
            if cells:
                lines.append("\t".join(cells))
    wb.close()
    return "\n".join(lines)


_MAX_UPLOAD_BYTES = 300 * 1024 * 1024  # 300 MB

# Extraction output caps. Tuned to comfortably fit the LLM's attention
# budget for a single KB-inference call while leaving headroom for the
# system prompt and existing-knowledge context.
#
# Rationale:
#   - 50k chars ≈ 25k tokens of input. claude-opus processes it in
#     ~15-25s and costs ~$0.12 at current prices; trivially fits any
#     modern context window.
#   - URLs tend to be single-page marketing copy — 10k is almost always
#     enough and Jina Reader has already stripped the chrome.
#   - Files (briefs, decks, whitepapers) are the variable-size source;
#     50k lets a 100-slide deck through after deduping, but caps
#     pathological uploads before they tank latency/cost/quality.
_EXTRACT_CAP_FILE = 50_000
_EXTRACT_CAP_URL = 10_000


def _normalize_extracted(text: str, max_chars: int) -> str:
    """Deduplicate exact-match lines and cap total length.

    Two low-cost cleanups that together kill the dominant sources of
    repetition in real-world source material:

    1. PPT decks repeat master slides (footer, page number, disclaimers,
       branding strap-lines) on every slide. Extracting every text frame
       verbatim produces N copies of the same 3-10 lines. Exact-match
       line dedup removes them while preserving the first occurrence so
       content order / `[Slide N]` markers stay intact.
    2. PDFs inherit the same footer/header repetition.
    3. URL fallback extracts nav/footer/cookie-banner links repeatedly
       from multi-column layouts — same pattern, same fix.

    We preserve single blank lines (they anchor section boundaries for
    the LLM), collapse runs of 3+ blank lines to a single break, and
    hard-cap at ``max_chars`` so pathologically large inputs can't
    blow up prompt length / LLM latency / cost.
    """
    seen: set[str] = set()
    lines: list[str] = []
    for line in text.splitlines():
        key = line.strip()
        if not key:
            # Keep blank lines — they carry paragraph structure. Collapse
            # runs later.
            lines.append("")
            continue
        if key in seen:
            continue
        seen.add(key)
        lines.append(line)

    out = "\n".join(lines)
    # Collapse 3+ consecutive newlines down to exactly 2 (one blank line).
    import re
    out = re.sub(r"\n{3,}", "\n\n", out).strip()
    return out[:max_chars]


async def extract_text_from_source(
    *,
    file: UploadFile | None,
    url: str | None,
    context_label: str = "extract-text",
) -> tuple[str, str, str | None]:
    """Shared helper: normalize a file-or-URL input into clean text.

    Both ``/ai/extract-text`` and ``/brandkits/{id}/extract-profile``
    consume the exact same set of input formats through the exact same
    pipeline — this helper is the single source of truth so they can't
    drift. (Historic drift caused PPTX to silently break on the brandkit
    path while working on the KB path.)

    Returns ``(normalized_text, source, filename)`` where ``source`` is
    ``"file"`` or ``"url"``. Raises ``HTTPException`` on unsupported
    format, oversized upload, empty result, or when neither input was
    provided.

    Applies ``_normalize_extracted`` (line dedup + length cap) before
    returning. ``context_label`` is only used in the log line.
    """
    if file and file.filename:
        declared = getattr(file, "size", None)
        if declared is not None and declared > _MAX_UPLOAD_BYTES:
            raise HTTPException(
                413,
                f"File too large: {declared / 1024 / 1024:.1f} MB. Maximum 300 MB.",
            )
        content = await file.read()
        if len(content) > _MAX_UPLOAD_BYTES:
            raise HTTPException(
                413,
                f"File too large: {len(content) / 1024 / 1024:.1f} MB. Maximum 300 MB.",
            )
        filename = file.filename.lower()

        if filename.endswith(".pdf"):
            text = _extract_pdf_text(content)
        elif filename.endswith((".docx", ".doc")):
            text = _extract_docx_text(content)
        elif filename.endswith((".xlsx", ".xls")):
            text = _extract_excel_text(content)
        elif filename.endswith((".pptx", ".ppt")):
            text = _extract_pptx_text(content)
        elif filename.endswith((".txt", ".csv")):
            text = content.decode("utf-8", errors="ignore")
        else:
            raise HTTPException(
                400,
                f"Unsupported file format: {file.filename}. Supported: PDF, Word, PowerPoint, Excel, CSV, TXT",
            )

        if not text.strip():
            raise HTTPException(400, "Failed to extract text content from the file")

        raw_len = len(text)
        normalized = _normalize_extracted(text, max_chars=_EXTRACT_CAP_FILE)
        logger.info(
            "%s | file=%s raw=%d normalized=%d (dedup+cap=%.1f%%)",
            context_label, file.filename, raw_len, len(normalized),
            100 * (1 - len(normalized) / raw_len) if raw_len else 0,
        )
        return normalized, "file", file.filename

    if url:
        if url.lower().endswith((".pdf", ".docx", ".doc")):
            try:
                async with httpx.AsyncClient(follow_redirects=True, timeout=15) as client:
                    resp = await client.get(url)
                    resp.raise_for_status()
            except Exception as e:
                raise HTTPException(400, f"Unable to access URL: {e}")

            if url.lower().endswith(".pdf"):
                text = _extract_pdf_text(resp.content)
            else:
                text = _extract_docx_text(resp.content)
        else:
            text = await _extract_url_text(url)

        if not text.strip():
            raise HTTPException(400, "Failed to extract text content from the URL")

        raw_len = len(text)
        normalized = _normalize_extracted(text, max_chars=_EXTRACT_CAP_URL)
        logger.info(
            "%s | url=%s raw=%d normalized=%d (dedup+cap=%.1f%%)",
            context_label, url[:80], raw_len, len(normalized),
            100 * (1 - len(normalized) / raw_len) if raw_len else 0,
        )
        return normalized, "url", None

    raise HTTPException(400, "Please provide a file or URL")


@router.post("/extract-text", response_model=ExtractTextResponse)
async def extract_text(
    file: UploadFile | None = File(None),
    url: str | None = Form(None),
):
    """Extract text from a PDF/Word/PowerPoint/Excel file or a URL."""
    text, source, filename = await extract_text_from_source(
        file=file, url=url, context_label="extract-text",
    )
    return ExtractTextResponse(text=text, source=source, filename=filename)


